import streamlit as st
from pdf2image import convert_from_bytes
from pptx import Presentation
from pptx.util import Inches, Pt
from io import BytesIO
from PIL import Image

# PDF를 이미지로 변환 후 PPTX 파일 생성 함수
def pdf_to_pptx(pdf_file):
    # PDF를 이미지로 변환
    images = convert_from_bytes(pdf_file.read())

    # PPTX 파일 생성
    presentation = Presentation()

    for image in images:
        # 슬라이드 크기를 이미지 크기에 맞게 설정
        width, height = image.size
        width_inches = width / 96  # 픽셀을 인치로 변환 (1인치 = 96px)
        height_inches = height / 96

        presentation.slide_width = Inches(width_inches)
        presentation.slide_height = Inches(height_inches)

        # 슬라이드 추가
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        image_stream = BytesIO()
        image.save(image_stream, format="PNG")

        # 이미지를 슬라이드에 추가
        slide.shapes.add_picture(image_stream, 0, 0, Inches(width_inches), Inches(height_inches))

    # PPTX 파일 저장
    pptx_stream = BytesIO()
    presentation.save(pptx_stream)
    pptx_stream.seek(0)
    return pptx_stream

# Streamlit 앱 설정
st.title("📄 PDF를 PPTX로 변환기 🖼️")
st.write("PDF 파일을 업로드하면 이미지로 변환 후 PPTX 파일을 생성합니다.")

# 파일 업로드
uploaded_file = st.file_uploader("PDF 파일을 업로드하세요", type=["pdf"])

if uploaded_file is not None:
    st.info("PDF 파일이 성공적으로 업로드되었습니다! 변환 중입니다...⏳")
    
    # PDF -> PPTX 변환
    try:
        pptx_file = pdf_to_pptx(uploaded_file)
        st.success("🎉 변환이 완료되었습니다! 이제 PPTX 파일을 다운로드하세요.")

        # 다운로드 링크 생성
        st.download_button(
            label="📥 PPTX 파일 다운로드",
            data=pptx_file,
            file_name="converted_presentation.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )
    except Exception as e:
        st.error(f"오류가 발생했습니다: {e}")
